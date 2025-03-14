import streamlit as st
import google.generativeai as genai
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings


class DocumentRAGChatbot:
    def __init__(self):
        # Initialize session state variables if they don't exist
        if "messages" not in st.session_state:
            st.session_state.messages = []
        if "vector_store" not in st.session_state:
            st.session_state.vector_store = None
        if "current_file" not in st.session_state:
            st.session_state.current_file = None

    def save_api_key(self, api_key):
        """Save and validate API key"""
        if api_key:
            genai.configure(api_key=api_key)
            st.session_state.api_key = api_key
            st.success("API key saved successfully!")
            return True
        else:
            st.error("Please enter an API key")
            return False

    def process_document(self, file):
        try:
            # Check if we already have processed this file
            if (st.session_state.vector_store is not None and 
                st.session_state.current_file == file.name):
                st.info(f"Document {file.name} is already processed and ready to use!")
                return
            
            # Get file extension
            file_extension = file.name.lower().split('.')[-1]
            
            # Map of supported file types to their processing functions
            processors = {
                'pdf': self.extract_text_from_pdf,
                'docx': self.extract_text_from_docx,
                'pptx': self.extract_text_from_pptx
            }
            
            # Check if file type is supported
            if file_extension not in processors:
                supported_formats = ', '.join(processors.keys())
                raise ValueError(f"Unsupported file format. Please upload one of: {supported_formats}")
            
            # Process document with appropriate function
            progress_text = st.empty()
            progress_bar = st.progress(0)
            cancel_button = st.button("Cancel Processing")
            
            try:
                # Update progress
                progress_text.text(f"Processing {file_extension.upper()} document...")
                progress_bar.progress(25)
                
                if cancel_button:
                    st.warning("Processing cancelled by user")
                    return
                
                texts = processors[file_extension](file)
                
                if not texts:
                    raise ValueError(f"No text content could be extracted from the {file_extension.upper()} file")
                
                # Update progress
                progress_text.text("Creating vector store...")
                progress_bar.progress(50)
                
                if cancel_button:
                    st.warning("Processing cancelled by user")
                    return
                
                # Update session state
                st.session_state.vector_store = self.create_vector_store(texts)
                st.session_state.current_file = file.name
                
                # Final progress update
                progress_text.text("Processing complete!")
                progress_bar.progress(100)
                st.success(f"Successfully processed {file.name}")
                
            finally:
                # Clean up progress indicators
                progress_text.empty()
                progress_bar.empty()
                
        except Exception as e:
            st.error(f"Error processing document: {str(e)}")
            st.session_state.vector_store = None
            st.session_state.current_file = None
            raise

    def extract_text_from_pdf(self, file):
        """Extract text from PDF files with robust error handling"""
        try:
            import PyPDF2
            import io
            
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file.getvalue()))
            texts = []
            
            # Validate PDF has pages
            if len(pdf_reader.pages) == 0:
                raise ValueError("The PDF file appears to be empty")
                
            # Extract text from each page with progress bar
            progress_bar = st.progress(0)
            for i, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                if text.strip():  # Only add non-empty pages
                    texts.append(text)
                progress_bar.progress((i + 1) / len(pdf_reader.pages))
            
            if not texts:
                raise ValueError("No readable text found in the PDF")
                
            return texts
            
        except Exception as e:
            st.error(f"Error extracting text from PDF: {str(e)}")
            raise

    def extract_text_from_docx(self, file):
        """Extract text from Word documents with robust error handling"""
        try:
            import docx
            import io
            
            doc = docx.Document(io.BytesIO(file.getvalue()))
            texts = []
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                if para.text.strip():  # Only add non-empty paragraphs
                    texts.append(para.text)
            
            if not texts:
                raise ValueError("No readable text found in the Word document")
                
            return texts
            
        except Exception as e:
            st.error(f"Error extracting text from Word document: {str(e)}")
            raise

    def extract_text_from_pptx(self, file):
        """Extract text from PowerPoint files with robust error handling"""
        try:
            from pptx import Presentation
            import io
            
            prs = Presentation(io.BytesIO(file.getvalue()))
            texts = []
            
            # Validate presentation has slides
            if len(prs.slides) == 0:
                raise ValueError("The PowerPoint file appears to be empty")
            
            # Extract text from each slide with progress bar
            progress_bar = st.progress(0)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():  # Only add non-empty shapes
                            slide_text.append(shape.text)
                if slide_text:  # Only add slides with text
                    texts.append("\n".join(slide_text))
                progress_bar.progress((i + 1) / len(prs.slides))
            
            if not texts:
                raise ValueError("No readable text found in the PowerPoint file")
                
            return texts
            
        except Exception as e:
            st.error(f"Error extracting text from PowerPoint: {str(e)}")
            raise

    def create_vector_store(self, texts):
        try:
            # Input validation
            if not texts or not isinstance(texts, list) or not all(isinstance(t, str) for t in texts):
                raise ValueError("Invalid input: texts must be a non-empty list of strings")

            # Create progress containers
            progress_text = st.empty()
            progress_bar = st.progress(0)
            
            try:
                # Configure text splitter
                progress_text.text("Configuring text splitter...")
                progress_bar.progress(10)
                
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=400,
                    length_function=len,
                    separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""],
                    keep_separator=True
                )

                # Process chunks
                progress_text.text("Creating document chunks...")
                progress_bar.progress(30)
                chunks = text_splitter.create_documents(texts)
                
                # Validate chunk creation
                if not chunks:
                    raise ValueError("No chunks were created from the input texts")
                
                progress_text.text(f"Created {len(chunks)} chunks from the document")
                progress_bar.progress(50)

                # Initialize embeddings
                progress_text.text("Initializing embeddings model...")
                progress_bar.progress(70)
                
                # Get API key from session state
                if "api_key" not in st.session_state:
                    raise ValueError("API key not found. Please enter your Google API key in the sidebar.")
                    
                try:
                    embeddings = GoogleGenerativeAIEmbeddings(
                        model="models/embedding-001",
                        google_api_key=st.session_state.api_key
                    )
                except Exception as e:
                    st.error("Failed to initialize embeddings model. Please check your API key.")
                    raise Exception(f"Embeddings initialization failed: {str(e)}")

                # Create vector store
                progress_text.text("Creating vector store...")
                progress_bar.progress(90)
                vector_store = FAISS.from_documents(chunks, embeddings)
                
                progress_text.text("Vector store created successfully!")
                progress_bar.progress(100)
                return vector_store

            finally:
                # Clean up progress indicators
                progress_text.empty()
                progress_bar.empty()

        except Exception as e:
            st.error(f"Error creating vector store: {str(e)}")
            raise

    def query_document(self, query: str, model_name: str = "gemini-1.5-flash") -> str:
        """Query the document content using RAG"""
        if "api_key" not in st.session_state:
            raise ValueError("API key not found. Please enter your Google API key in the sidebar.")
            
        results = st.session_state.vector_store.similarity_search(query, k=5)
        context = "\n".join([doc.page_content for doc in results])

        model = genai.GenerativeModel(model_name)
        prompt = f"""You are a knowledgeable and precise AI assistant. Your task is to provide accurate, well-structured responses based on the given context.

Instructions:
1. Analyze the provided context thoroughly
2. Answer the question using ONLY information from the context
3. Structure your response in a clear, organized manner
4. If relevant, use bullet points or numbered lists for better readability
5. Include specific details and examples from the context to support your answer
6. If the context doesn't contain enough information, clearly state: "I cannot provide a complete answer as the context doesn't contain sufficient information about [specific aspect]"
7. If the question is unclear, ask for clarification
8. Maintain a professional yet conversational tone
9. Easy english for faster understanding

Context:
{context}

Question:
{query}

Response:
"""
        response = model.generate_content(prompt)
        return response.text


    def generate_document_summary(self) -> str:
        """Generate a comprehensive summary of the document"""
        try:
            if not st.session_state.vector_store:
                raise ValueError("No document has been processed yet")

            prompt = """Generate a comprehensive summary of the document. The summary should:
            1. Be well-structured and easy to read
            2. Cover the main topics and themes
            3. Be around 3-4 paragraphs long
            4. Maintain the key message and intent of the original document
            
            Document content:
            {context}
            
            Summary:"""

            # Get all document chunks for complete context
            results = st.session_state.vector_store.similarity_search(
                "what is this document about?", k=10
            )
            context = "\n".join([doc.page_content for doc in results])

            model = genai.GenerativeModel("gemini-1.5-flash")
            response = model.generate_content(prompt.format(context=context))
            return response.text

        except Exception as e:
            st.error(f"Error generating summary: {str(e)}")
            return None

    def extract_key_points(self) -> str:
        """Extract key points from the document"""
        try:
            if not st.session_state.vector_store:
                raise ValueError("No document has been processed yet")

            prompt = """Extract the most important key points from the document. Your response should:
            1. List 5-7 key points in bullet form
            2. Each point should be clear and concise
            3. Focus on the most significant information
            4. Maintain factual accuracy
            5. Order points by importance
            
            Document content:
            {context}
            
            Key Points:"""

            # Get document chunks for context
            results = st.session_state.vector_store.similarity_search(
                "what are the main points of this document?", k=10
            )
            context = "\n".join([doc.page_content for doc in results])

            model = genai.GenerativeModel("gemini-1.5-flash")
            response = model.generate_content(prompt.format(context=context))
            return response.text

        except Exception as e:
            st.error(f"Error extracting key points: {str(e)}")
            return None


def main():
    st.set_page_config(page_title="DocsMind", layout="wide")

    st.title("DocsMind - Document Analysis & Chat")
    st.write(
        "DocsMind is a powerful document-based chatbot that uses Retrieval-Augmented Generation (RAG) to provide accurate answers from your documents using Google's Generative AI."
    )
    
    # Initialize chatbot
    chatbot = DocumentRAGChatbot()

    # API key input
    api_key = st.sidebar.text_input("Enter your Google API key:", type="password")
    if api_key:
        if chatbot.save_api_key(api_key):
            st.sidebar.success("API key saved successfully!")
        else:
            st.sidebar.error("Invalid API key")

    # File upload
    uploaded_file = st.sidebar.file_uploader(
        "Upload a document (PDF, DOCX, or PPTX)", type=["pdf", "docx", "pptx"]
    )

    if uploaded_file:
        # Display file information
        st.subheader("Document Information")
        col1, col2, col3 = st.columns(3)
        col1.metric("File Name", uploaded_file.name)
        col2.metric("File Size", f"{uploaded_file.size / 1024:.1f} KB")
        col3.metric("File Type", uploaded_file.name.split('.')[-1].upper())

        # Add process button
        if st.button("Process Document"):
            # Only process if it's a new file or no vector store exists
            if (st.session_state.get("current_file") != uploaded_file.name or 
                st.session_state.get("vector_store") is None):
                chatbot.process_document(uploaded_file)

    # Initialize session state for expanders if not exists
    if "show_summary" not in st.session_state:
        st.session_state.show_summary = False
    if "show_keypoints" not in st.session_state:
        st.session_state.show_keypoints = False
    if "summary_content" not in st.session_state:
        st.session_state.summary_content = None
    if "keypoints_content" not in st.session_state:
        st.session_state.keypoints_content = None

    # Create two columns for the main content with 2:3 ratio
    col1, col2 = st.columns([2, 3])

    # Document analysis tools in the first column (smaller)
    with col1:
        st.subheader("Document Analysis")
        if st.session_state.get("vector_store"):
            # Summary section
            with st.expander("📝 Document Summary", expanded=st.session_state.get("show_summary", False)):
                if st.button("Generate Summary", key="summary_btn", use_container_width=True):
                    with st.spinner("Generating summary..."):
                        summary = chatbot.generate_document_summary()
                        if summary:
                            st.session_state["summary_content"] = summary
                            st.session_state["show_summary"] = True
                
                if st.session_state.get("summary_content"):
                    st.markdown(st.session_state["summary_content"])
            
            # Key points section
            with st.expander("🔑 Key Points", expanded=st.session_state.get("show_keypoints", False)):
                if st.button("Extract Key Points", key="keypoints_btn", use_container_width=True):
                    with st.spinner("Extracting key points..."):
                        keypoints = chatbot.extract_key_points()
                        if keypoints:
                            st.session_state["keypoints_content"] = keypoints
                            st.session_state["show_keypoints"] = True
                
                if st.session_state.get("keypoints_content"):
                    st.markdown(st.session_state["keypoints_content"])

    # Chat interface in the second column (larger)
    with col2:
        st.subheader("Chat with Document")
        if st.session_state.get("vector_store"):
            # Display chat messages
            for message in st.session_state.get("messages", []):
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])

            # Chat input
            if prompt := st.chat_input("Enter your message"):
                # Add user message to chat history
                st.session_state.messages.append({"role": "user", "content": prompt})
                
                # Display user message
                with st.chat_message("user"):
                    st.markdown(prompt)

                # Generate and display assistant response
                with st.chat_message("assistant"):
                    with st.spinner("Thinking..."):
                        response = chatbot.query_document(prompt)
                        st.markdown(response)
                        st.session_state.messages.append({"role": "assistant", "content": response})
        else:
            st.info("Please upload a document to start chatting!") 
    st.write("Made by [Naman0807](https://github.com/Naman0807) can contribute at [here](https://github.com/Naman0807/DocsMind)") 
if __name__ == "__main__":
    main()
